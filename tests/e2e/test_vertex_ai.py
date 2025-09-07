import os
import pytest
import pathlib
import json
import uuid
import logging
import time

# Configure logging to see test progress
logging.basicConfig(level=logging.INFO)
log = logging.getLogger(__name__)

# --- Step 1: Import all required libraries ---
from PIL import Image, ImageDraw
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import docx
import pptx
import imageio
import numpy as np
from pydub import AudioSegment
from pydub.generators import Sine
from google.cloud import storage
from google.api_core import exceptions
import vertexai
from vertexai.generative_models import GenerativeModel, Part


# --- Step 2: Define Fixtures for File Generation and GCS Management ---

@pytest.fixture(scope="session")
def gcs_client():
    """
    Initializes the Google Cloud Storage client.
    It automatically uses Application Default Credentials (from gcloud CLI).
    """
    
    
    try:
        # The client automatically discovers the project_id from the environment.
        return storage.Client()
    except Exception as e:
        pytest.fail(f"Could not initialize GCS Client. Ensure you are authenticated with 'gcloud auth application-default login'. Error: {e}")

@pytest.fixture(scope="session")
def test_bucket(gcs_client):
    """
    Ensures the hardcoded test bucket exists for the session.
    It does NOT create or delete the bucket, assuming it's permanent.
    """
    bucket_name = "vertex-ai-test-files"
    bucket = gcs_client.bucket(bucket_name)
    if not bucket.exists():
        pytest.fail(f"The required GCS bucket '{bucket_name}' does not exist. Please create it manually.")
    
    log.info(f"Using existing GCS test bucket: {bucket_name}")
    yield bucket

    # Teardown: Clean up blobs created during the test run
    log.info(f"Cleaning up blobs in GCS test bucket: {bucket_name}")
    blobs_to_delete = []
    # Only blobs with the 'test-run-' prefix are deleted for safety
    for blob in bucket.list_blobs(prefix="test-run-"):
        blobs_to_delete.append(blob.name)
    
    if blobs_to_delete:
        try:
            # Use batch deletion for efficiency
            bucket.delete_blobs(blobs_to_delete)
            log.info(f"Deleted {len(blobs_to_delete)} test blobs.")
        except Exception as e:
            log.error(f"Failed to clean up all blobs in {bucket_name}. You may need to delete them manually. Error: {e}")

@pytest.fixture(scope="session")
def file_generator():
    """
    Returns a function that generates ALL file types in a project-local directory.
    """
    
    output_dir = pathlib.Path("generated_test_files")
    output_dir.mkdir(exist_ok=True)

    def _generator(mime_type: str):
        # Create a unique subdirectory for each file to avoid name collisions
        # and keep the output organized, mimicking the original safe behavior.
        file_specific_dir = output_dir / f"run-{uuid.uuid4().hex[:6]}"
        file_specific_dir.mkdir()
        
        temp_dir = file_specific_dir
        # Handles complex mime types to extract a valid extension
        extension = mime_type.split("/")[-1].split(".")[-1].split("-")[-1]
        file_path = temp_dir / f"sample.{extension}"

        # --- Generation logic for EACH file type ---
        if mime_type.startswith("image/"):
            img = Image.new('RGB', (120, 80), color='blue')
            d = ImageDraw.Draw(img)
            d.text((10, 10), f"{extension.upper()} Test", fill=(255, 255, 0))
            # HEIC/HEIF are not standard output formats, so we save as PNG.
            # The mime_type in the API call is what matters.
            if extension in ['heic', 'heif']:
                file_path = file_path.with_suffix('.png')
            img.save(file_path)
        elif mime_type.startswith("video/"):
            # Create a simple video with a few frames of solid color
            frames = []
            for _ in range(10):  # 10 frames
                frame = np.full((128, 128, 3), [255, 0, 0], dtype=np.uint8)
                frames.append(frame)
            imageio.mimwrite(file_path, frames, fps=10, format=extension)
        elif mime_type.startswith("audio/"):
            sine_wave = Sine(440).to_audio_segment(duration=1000)
            if extension == "aac":
                sine_wave.export(file_path, format="adts", codec="aac")
            else:
                sine_wave.export(file_path, format=extension)
        elif mime_type == "application/pdf":
            c = canvas.Canvas(str(file_path), pagesize=letter)
            c.drawString(72, 800, "This is a test PDF document.")
            c.save()
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document()
            doc.add_paragraph("Test paragraph in a DOCX file.")
            doc.save(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = pptx.Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "PPTX Test Slide"
            prs.save(file_path)
        elif mime_type == "application/rtf":
            # Generate a simple RTF
            file_path.write_text(r"{\rtf1\ansi\deff0 {\fonttbl{\f0 Arial;}}\f0\fs24 RTF content test.}")
        elif mime_type == "text/html":
            file_path.write_text("<!DOCTYPE html><html><head><title>Test Title</title></head><body><p>Test paragraph.</p></body></html>")
        elif mime_type == "text/csv":
            file_path.write_text("id,name\n1,csv_test")
        elif mime_type == "application/json":
            file_path.write_text(json.dumps({"id": 1, "object": "json_test"}))
        elif mime_type == "application/jsonl":
            file_path.write_text('{"id": 1, "object": "line1"}\n{"id": 2, "object": "line2"}')
        elif mime_type in ["application/xml", "text/xml"]:
            file_path.write_text("<root><id>1</id><object>xml_test</object></root>")
        else:  # Covers text/plain and others
            file_path.write_text("This is a plain text file for testing.")

        log.info(f"Generated test file: {file_path}")
        return file_path

    return _generator


@pytest.fixture(scope="module")
def vertex_ai_model():
    """
    Initializes the Vertex AI Generative Model for the test module.
    It automatically discovers the project_id from the environment (gcloud CLI).
    """
    
    
    try:
        # project and location are discovered from the environment automatically
        vertexai.init()
        model = GenerativeModel("gemini-2.5-pro")
        return model
    except Exception as e:
        pytest.fail(f"Failed to initialize Vertex AI Model. Error: {e}")

# --- Step 3: Define Test Cases and the Main Test Function ---

# The most comprehensive list of file types supported by Vertex AI Gemini
ALL_SUPPORTED_FORMATS = [
    # Images
    ("image/png", "Describe the color of the background and the text."),
    ("image/jpeg", "Describe the color of the background and the text."),
    ("image/webp", "Describe the color of the background and the text."),
    ("image/heic", "Describe the color of the background and the text."),
    ("image/heif", "Describe the color of the background and the text."),
    # Videos
    ("video/mp4", "What is the main color in this short video?"),
    ("video/mov", "What is the main color in this short video?"),
    ("video/mpeg", "What is the main color in this short video?"),
    ("video/avi", "What is the main color in this short video?"),
    ("video/wmv", "What is the main color in this short video?"),
    # Audio
    ("audio/mp3", "Does this audio contain a sine wave tone?"),
    ("audio/wav", "Does this audio contain a sine wave tone?"),
    ("audio/aac", "Does this audio contain a sine wave tone?"),
    ("audio/ogg", "Does this audio contain a sine wave tone?"),
    ("audio/flac", "Does this audio contain a sine wave tone?"),
    # Documents
    ("application/pdf", "What is the subject of this PDF document?"),
    ("text/html", "What is the title of this HTML document?"),
    ("text/csv", "What is the name in the second row of this CSV?"),
    ("text/plain", "What is this text file about?"),
]

@pytest.mark.parametrize("mime_type, prompt", ALL_SUPPORTED_FORMATS)
def test_real_analysis_of_all_formats(vertex_ai_model, test_bucket, file_generator, mime_type, prompt):
    """
    Performs a REAL end-to-end test for every supported file format.
    It generates a file, uploads it to a GCS bucket, and sends it to the Vertex AI API.
    """
    log.info(f"--- Starting test for MIME type: {mime_type} ---")
    
    local_file_path = file_generator(mime_type)
    assert local_file_path.exists(), "File generator failed to create a file."

    blob_name = f"test-run-{uuid.uuid4().hex}/{local_file_path.name}"
    blob = test_bucket.blob(blob_name)
    try:
        log.info(f"Uploading {local_file_path} to gs://{test_bucket.name}/{blob.name}")
        blob.upload_from_filename(str(local_file_path))
    except Exception as e:
        pytest.fail(f"Failed to upload file to GCS. Error: {e}")

    gcs_uri = f"gs://{test_bucket.name}/{blob.name}"
    file_part = Part.from_uri(uri=gcs_uri, mime_type=mime_type)

    response_text = ""
    try:
        log.info(f"Sending prompt to Vertex AI for {gcs_uri}")
        response = vertex_ai_model.generate_content([prompt, file_part])
        response_text = response.text
        log.info(f"Received response from Vertex AI: '{response_text[:100].strip()}...'")
    except Exception as e:
        pytest.fail(f"Vertex AI API call failed for {mime_type}. Error: {e}")

    assert response, f"Received an empty response for {mime_type}."
    assert response_text, f"Response text is empty for {mime_type}."
    if "error" in response_text.lower():
        log.error(f"Response for {mime_type} contains an error message: {response_text}")
        assert "error" not in response_text.lower()
    
    # Simple content checks
    if mime_type.startswith("image"):
        assert "blue" in response_text.lower() and ("yellow" in response_text.lower() or "white" in response_text.lower() or "black" in response_text.lower())
    if mime_type.startswith("video"):
         assert "red" in response_text.lower()
    if mime_type.startswith("audio"):
         assert "sine wave" in response_text.lower() or "tone" in response_text.lower()
    if "json" in mime_type:
        assert "json_test" in response_text.lower() or "line2" in response_text.lower()
    if "xml" in mime_type:
        assert "xml_test" in response_text.lower()

    log.info(f"--- Test PASSED for MIME type: {mime_type} ---")
    
    # API rate limiting to avoid "429 Resource exhausted" errors
    time.sleep(5) # Increased to 5 seconds for more safety
